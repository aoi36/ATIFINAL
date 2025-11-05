# scraper_service.py
import os
import time
import json
import re
import sqlite3
import requests
import urllib.parse
from selenium import webdriver
from datetime import timedelta
import docx
import pptx
import pdfplumber
from selenium.webdriver.support.ui import WebDriverWait
from selenium.common.exceptions import TimeoutException
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.by import By
# ... (all other imports needed for scraping: requests, bs4, docx, pptx, pdfplumber, etc.)
from datetime import datetime
from dateutil import parser as date_parser

# Import from our new modules
import state
from config import (
    DATABASE_FILE, SAVE_DIR, LMS_USERNAME, LMS_PASSWORD, STATE_FILE,
    REQUESTS_TIMEOUT, MAX_SUBPAGES, MAX_TEXT_LENGTH_FOR_SUMMARY
)
from search_service import clear_search_index
# Note: AI functions are no longer called from here, so we don't import them.
# Import the file-reading and deadline-parsing helpers
# (Paste clean_file_text, parse_time_remaining, read_docx, read_pptx, read_pdf, download_file here)
#
# Example (paste your full function):

def download_file(url, folder, cookies, headers, link_text="") -> str | None:
    filename = None
    local_path = None
    print(f"         [Download Func] Attempting HEAD for: {url}") # DEBUG
    try:
        final_url = url
        content_disposition = None
        # Use HEAD first, follow redirect once if needed
        # Disable redirects temporarily to inspect original headers
        with requests.head(url, headers=headers, cookies=cookies, timeout=REQUESTS_TIMEOUT, allow_redirects=False) as r_head:
            print(f"         [Download Func] HEAD status: {r_head.status_code}") # DEBUG
            # If redirected, follow it once to get the final URL/headers
            if 300 <= r_head.status_code < 400 and 'Location' in r_head.headers:
                final_url = urllib.parse.urljoin(url, r_head.headers['Location'])
                print(f"         [Download Func] Redirect detected. Following to: {final_url}") # DEBUG
                # Make HEAD request to the final URL
                with requests.head(final_url, headers=headers, cookies=cookies, timeout=REQUESTS_TIMEOUT) as r_final_head:
                    print(f"         [Download Func] Final HEAD status: {r_final_head.status_code}") # DEBUG
                    content_disposition = r_final_head.headers.get('content-disposition')
                    # Raise error for the final URL if it fails
                    r_final_head.raise_for_status()
            else:
                # Use headers from the original URL if no redirect
                content_disposition = r_head.headers.get('content-disposition')
                # Raise error for the original URL if it failed
                r_head.raise_for_status()


        url = final_url # Use final URL for the GET request
        print(f"         [Download Func] Using final URL for download: {url}") # DEBUG


        # --- Filename Determination ---
        filename_source = "Unknown" # DEBUG
        if content_disposition:
            # More robust regex to handle filename*=UTF-8''... encoding
            disp_match = re.search(r'filename\*?=(?:UTF-\d{1,2}\'\')?([^";\n]+)', content_disposition, re.IGNORECASE)
            if disp_match:
                filename = urllib.parse.unquote(disp_match.group(1)).strip('"')
                filename_source = "Content-Disposition Header" # DEBUG
                print(f"         [Download Func] Filename from header: {filename}")


        # Fallback using URL path
        if not filename:
             parsed_url = urllib.parse.urlparse(url)
             filename = os.path.basename(parsed_url.path)
             if filename:
                 filename_source = "URL Path" # DEBUG
                 print(f"         [Download Func] Filename from URL path: {filename}")


        # Fallback using link text if other methods fail or produce generic names
        # Added empty string check for filename
        if not filename or filename in ["pluginfile.php", "download", "content", "index.php", ""]:
             if link_text and any(link_text.lower().endswith(ext) for ext in (".pdf", ".docx", ".pptx", ".zip", ".rar", ".xls", ".xlsx", ".txt", ".csv")):
                  filename = re.sub(r'[\\/*?:"<>|]', "_", link_text).strip()
                  filename_source = "Link Text" # DEBUG
                  print(f"         [Download Func] Filename inferred from link text: {filename}")
             else:
                  # Last resort, generate a generic name
                  filename_source = "Generic Fallback" # DEBUG
                  generic_name = "downloaded_file"
                  query = urllib.parse.urlparse(url).query.replace('&','_').replace('=','-')
                  if query:
                       safe_query = re.sub(r'[\\/*?:"<>|]', "_", query)[:50] # Limit length
                       generic_name += "_" + safe_query
                  else:
                       # Add timestamp if no query to help uniqueness
                       generic_name += f"_{int(time.time())}"


                  # Try to guess extension based on link text or URL
                  ext_match = re.search(r'\.(pdf|docx|pptx|zip|rar|xls|xlsx|txt|csv)$', link_text or url, re.IGNORECASE)
                  if ext_match:
                      generic_name += ext_match.group(0)


                  filename = generic_name
                  print(f"         [Download Func] Using generic filename: {filename}")


        # Final sanitization and path creation
        filename = re.sub(r'[\\/*?:"<>|\n\r\t]', "_", filename).strip()[:200] # Shorten long names too
        local_path = os.path.join(folder, filename)
        print(f"         [Download Func] Determined path (from {filename_source}): {local_path}") # DEBUG


        # Perform the actual download using GET
        print(f"         ⬇️ Starting GET request for download...") # DEBUG
        with requests.get(url, headers=headers, cookies=cookies, timeout=REQUESTS_TIMEOUT*2, stream=True) as r_get: # Longer timeout for download
            print(f"         [Download Func] GET request status: {r_get.status_code}") # DEBUG
            r_get.raise_for_status() # Check for HTTP errors (4xx, 5xx)
            print(f"         [Download Func] Writing to file...") # DEBUG
            bytes_written = 0
            with open(local_path, "wb") as f:
                for chunk in r_get.iter_content(chunk_size=8192*4): # Slightly larger chunk
                    if chunk: # filter out keep-alive new chunks
                        f.write(chunk)
                        bytes_written += len(chunk)
            print(f"         [Download Func] Finished writing {bytes_written} bytes.") # DEBUG


        # Check if file is empty, might indicate an issue despite 200 OK
        if bytes_written == 0 and os.path.exists(local_path):
             print(f"         ⚠️ WARNING: Downloaded file is empty: {local_path}")
             # Optionally delete empty file: os.remove(local_path); return None


        print(f"         📥 Download complete -> {local_path}")
        return local_path


    # --- Specific Exception Handling ---
    except requests.exceptions.HTTPError as http_e:
        print(f"         ⚠️ Failed download {url} (HTTP Error): {http_e.response.status_code} {http_e.response.reason}") # DEBUG
        return None
    except requests.exceptions.Timeout:
        print(f"         ⚠️ Failed download {url} (Timeout after {REQUESTS_TIMEOUT}s or {REQUESTS_TIMEOUT*2}s)") # DEBUG
        return None
    except requests.exceptions.ConnectionError as conn_e:
        print(f"         ⚠️ Failed download {url} (Connection Error): {conn_e}") # DEBUG
        return None
    except requests.exceptions.RequestException as req_e:
        print(f"         ⚠️ Failed download {url} (General Request Error): {req_e}") # DEBUG
        return None
    except Exception as e:
        print(f"         ⚠️ Failed download {url} (Other Error): {e}") # DEBUG
        # Clean up partial file if it exists and path was determined
        if local_path and os.path.exists(local_path):
             try:
                 # Check size before deleting, maybe small files are valid
                 if os.path.getsize(local_path) == 0:
                     os.remove(local_path)
                     print(f"         🧹 Cleaned up empty/partial file: {local_path}")
             except OSError: pass
        return None


def clean_file_text(text: str) -> str:
    """
Recap of the Whoosh Search Implementation
Here's a breakdown of the Whoosh integration we've built, explaining what each part does:

## 1. Imports and Configuration (app.py)

At the top of your app.py, we added:

* whoosh.index: Manages the index files (create, open, check existence).
* whoosh.fields: Defines the "schema" (the "columns" of your search database).
* whoosh.qparser: Parses user query strings (like "Java TCP") into a format Whoosh understands.
* whoosh.highlight: Contains classes to format the search results, like Formatter (which you correctly identified as the base class) and ContextFragmenter (which we use to create snippets).
* INDEX_DIR = "search_index": A configuration variable telling the script where to store the index files.

## 2. Whoosh Helper Functions

We added a few helper functions to manage the index:

* get_search_schema(): This is the blueprint for your search index.
    * course_id, course_name, file_name, file_type: These are set as STORED, meaning you can get this data back in your search results.
    * content: This is set as TEXT(stored=True, phrase=True).
        * TEXT: This is the crucial part. It tells Whoosh to analyze this text, break it down into searchable words (tokens), and index them.
        * stored=True: Saves a copy of the text within the index, which allows Whoosh to generate snippets (like in your JSON output).
        * phrase=True: Allows for phrase searching (e.g., searching for "Java Socket" as an exact phrase).
* get_or_create_index(): A utility function to safely open the index if it exists or create it if it doesn't.
* clear_search_index(): Correctly re-creates the index from scratch, which is necessary before running a full scrape to remove old/deleted file entries.
* SimpleFormatter(Formatter): A small custom class to wrap the highlighted search terms in <strong> tags, as seen in your JSON snippet.

## 3. Modifying perform_full_scrape (The Scraper)

This is where the index is built:

* Clear & Open Writer: At the *start* of the scrape (after fetching the course list), it calls clear_search_index() to get a fresh, empty index and opens a single index_writer.
* Add Documents: As the script scrapes each course and its files (both direct and nested), after successfully extracting text from a file (e.g., with read_pptx), it does this:
    * print(f"[Search] Indexing {filename}...")
    * index_writer.add_document(course_id=..., content=extracted_text, ...)
    * This adds the file's text and metadata to the writer.
* Commit Changes: At the very *end* of the entire scrape (after the main course loop finishes), it calls index_writer.commit(). This is a single, efficient operation that saves all the new documents to the index on disk.

## 4. The /api/search Endpoint (The Search Engine)

This is what Postman is calling:

* Get Query: It gets the search term from the URL parameter (e.g., ?q=TCP).
* Open Index: It opens the existing index in read-only mode (using ix = open_dir(INDEX_DIR)).
* Check for Empty: It correctly checks if ix.doc_count() == 0.
* Use Searcher: It opens an ix.searcher() to perform the query.
* Parse Query: It uses QueryParser("content", ...) to tell Whoosh, "Search for the user's terms primarily in the content field."
* Search: It calls searcher.search(query, limit=10).
* Format Results:
    * results.formatter = SimpleFormatter(): Tells Whoosh to use our custom <strong> tags for highlighting.
    * results.fragmenter = ContextFragmenter(maxchars=200, surround=50): This is what generates the "snippet". It finds the search term(s) in the stored content, grabs about 50 characters before and after (up to 200 total), and wraps the terms in the formatter.
* Build JSON: It loops through the results (the "hits") and builds the JSON list you see, pulling the stored fields (course_id, file_name) and the newly generated snippet.
* Return: It sends the JSON list back.

The output you're seeing, with the "snippet" containing HTML-like <strong> tags, is the *intended and correct* behavior of this search setup. The frontend (like a web browser) will interpret <strong>...</strong> as "make this text bold," highlighting the search term for the user.
    Removes common slide/text garbage like '---' and unicode artifacts."""
    if not text:
        return ""
    # Remove PowerPoint/PDF slide separators
    text = re.sub(r'\n---\n', '\n', text) # Remove '---' separators on their own lines
    # Remove unicode artifacts like \u000b (vertical tab)
    text = text.replace('\u000b', ' ')
    # Replace multiple newlines/spaces with a single space
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def parse_time_remaining(time_str: str) -> timedelta | None:
    """
    Parses Vietnamese 'Time Remaining' strings like 'Còn lại X Các ngày Y giờ'.
    Returns a timedelta object representing the duration, or None if parsing fails.
    """
    days = 0
    hours = 0
    # Use regex to find days and hours, making numbers optional
    # Handles "Còn lại X Các ngày Y giờ", "X Các ngày", "Y giờ"
    match = re.search(r'(?:(\d+)\s*Các ngày)?\s*(?:(\d+)\s*giờ)?', time_str, re.IGNORECASE)


    if match:
        days_str, hours_str = match.groups()
        try:
            days = int(days_str) if days_str else 0
            hours = int(hours_str) if hours_str else 0
            # Basic validation: ensure at least one part was found and positive
            if days >= 0 and hours >= 0 and (days > 0 or hours > 0):
                 return timedelta(days=days, hours=hours)
            else:
                 # Handle cases like "0 giờ" or just "Còn lại" if regex matches unexpectedly
                 print(f"         [Parse Time] Parsed zero or negative duration from '{time_str}'")
                 return None
        except (ValueError, TypeError):
             print(f"         [Parse Time] Error converting parsed numbers from '{time_str}'")
             return None # Error converting days/hours to int
    else:
        # Add handling for minutes if needed, e.g., 'Z phút'
        match_minutes = re.search(r'(\d+)\s*phút', time_str, re.IGNORECASE)
        if match_minutes:
            try:
                minutes = int(match_minutes.group(1))
                if minutes >= 0:
                    # Approximation: If only minutes left, treat as due very soon
                    return timedelta(minutes=minutes)
            except (ValueError, TypeError):
                 print(f"         [Parse Time] Error converting parsed minutes from '{time_str}'")
                 return None


        print(f"         [Parse Time] Could not parse duration from '{time_str}'")
        return None # Pattern not found

def read_docx(file_path: str) -> str:
    # ... (Paste your full function code here) ...
    try: doc = docx.Document(file_path); return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception as e: print(f" [DOCX Error] {os.path.basename(file_path)}: {e}"); return ""

def read_txt(file_path: str) -> str:
    """Extracts text from a plain .txt file, trying common encodings."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        try: # Fallback for other common encodings
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception as e:
            print(f"         [TXT Error] Failed to read {os.path.basename(file_path)}: {e}")
            return ""
    except Exception as e:
        print(f"         [TXT Error] Failed to read {os.path.basename(file_path)}: {e}")
        return ""


def read_pptx(file_path: str) -> str:
    # ... (Paste your full function code here) ...
    try:
        prs = pptx.Presentation(file_path); full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                para_text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                if para_text: full_text.append(para_text)
        return "\n\n---\n\n".join(full_text)
    except Exception as e: print(f" [PPTX Error] {os.path.basename(file_path)}: {e}"); return ""


# --- [MODIFIED] File Reading Helper ---
def read_pdf(file_path: str) -> str:
    """Extracts all text from a .pdf file using pdfplumber."""
    full_text = []
    try:
        with pdfplumber.open(file_path) as pdf:
            # --- [FIX] ---
            # Removed the 'if pdf.is_encrypted:' block, as that attribute
            # does not exist for pdfplumber.
            # We'll catch password errors in the 'except' block instead.
            # --- [END FIX] ---

            for i, page in enumerate(pdf.pages):
                # extract_text can return None if page has no text
                text = page.extract_text(x_tolerance=1, y_tolerance=1)
                if text:
                    full_text.append(f"--- Page {i+1} ---\n{text}")
        return "\n\n".join(full_text)
    
    except pdfplumber.exceptions.PasswordRequired:
         # Add a specific catch for password-protected files
         print(f"         [PDF Error] Failed to read {os.path.basename(file_path)}: File is password-protected.")
         return ""
    except Exception as e:
        # Catch other errors (e.g., corrupted file)
        print(f"         [PDF Error] Failed to read {os.path.basename(file_path)}: {e}")
        return ""


# (Paste your AI-based deadline extractors here, as they are part of scraping)
from ai_service import ai_client # Need the client
from bs4 import BeautifulSoup
def extract_deadline_with_selectors(soup: BeautifulSoup) -> dict | None:
    """
    Tries to get deadline data using specific selectors and basic text pattern matching.
    """
    try:
        # 1. Try Specific Table Cells (Most Reliable if Present)
        selectors = [
            "td.submissionstatustable_duedate", # Specific class for due date cell?
            "td.submissionstatustable_timeremaining", # Specific class for time remaining?
            ".submissionstatustable td.c2", # Value cell in submission status table
            "td.timeremaining", # General time remaining cells
            "td.overdue",       # General overdue cells
            ".timeremaining",   # Div/Span with class timeremaining
            ".overdue"          # Div/Span with class overdue
            ]
        for sel in selectors:
            element = soup.select_one(sel)
            if element:
                text = element.get_text(strip=True)
                if not text: continue # Skip empty elements


                # Determine status based on selector or content
                status = "Unknown"
                if "duedate" in sel or "Due date" in element.find_previous_sibling(class_="c1").get_text() if element.find_previous_sibling(class_="c1") else False:
                     status = "Due"
                elif "overdue" in sel or "overdue" in text.lower() or "quá hạn" in text.lower():
                     status = "Overdue"
                elif "timeremaining" in sel or "remaining" in text.lower() or "còn lại" in text.lower():
                     status = "Time Remaining"
                elif re.search(r'\d{1,2}\s+\w+\s+\d{4},\s+\d{1,2}:\d{2}\s+(?:AM|PM)', text): # Basic check for absolute date format
                     status = "Due" # Assume absolute dates mean 'Due'


                if status != "Unknown":
                    print(f"         [Selector] Found deadline via selector '{sel}': {status} - {text}")
                    return {"status": status, "time": text}


        # 2. Look in Activity Introduction/Description Areas
        # These selectors might need adjustment based on your Moodle theme
        intro_selectors = ["#intro.box.generalbox", ".activity-description", ".no-overflow"]
        for sel in intro_selectors:
            intro_box = soup.select_one(sel)
            if intro_box:
                intro_text = intro_box.get_text(" ", strip=True) # Get text content
                # Use Regex to find potential deadlines (Vietnamese/English keywords)
                # Look for patterns like "Deadline:", "Hạn nộp:", followed by a date/time
                # Example: (\b(?:Deadline|Hạn nộp)\b[:\s]*)             # Keyword + separator
                #          (.*?(?:\d{1,2}/\d{1,2}/\d{4}|\w+\s+\d{1,2})) # Date part (dd/mm/yyyy or Month dd)
                #          (?:.*?(\d{1,2}:\d{2}\s*(?:AM|PM)?))?        # Optional time part
                deadline_pattern = r'(?i)(\b(?:Deadline|Hạn nộp|Due date|Hết hạn)\b[:\s]*)(.*?((?:Mon|Tue|Wed|Thu|Fri|Sat|Sun|Thứ)\s*\d+.*?\d{4}.*?\d{1,2}:\d{2}\s*(?:AM|PM)?|\d{1,2}[/-]\d{1,2}[/-]\d{2,4}.*?\d{1,2}:\d{2}\s*(?:AM|PM)?|\d+\s+(?:days|ngày|hours|giờ|minutes|phút)\b))'
                match = re.search(deadline_pattern, intro_text)
                if match:
                    status_keyword = match.group(1).strip(": ")
                    time_str = match.group(2).strip()
                    status = "Due" # Assume found date patterns mean 'Due'
                    # Refine status based on keyword if needed (though Due is likely)
                    if "Hạn nộp" in status_keyword or "Due" in status_keyword: status = "Due"


                    print(f"         [Selector/Regex] Found deadline via intro text pattern: {status} - {time_str}")
                    return {"status": status, "time": time_str}


        # 3. If nothing found by selectors or regex
        return None


    except Exception as e:
        print(f"         [Selector WARN] Error during advanced selector search: {e}")
        return None

def extract_deadline_with_ai(html_content: str) -> dict | None:
    # ... (Paste your full function code with rate limit handling here) ...
    if not ai_client: return None
    print("         [AI Fallback] Trying Gemini for deadline...")
    try:
        soup_for_ai = BeautifulSoup(html_content, 'html.parser')
        main_content = soup_for_ai.select_one("#region-main")
        html_to_send = str(main_content) if main_content else html_content
    except Exception: html_to_send = html_content
    if len(html_to_send) > 100000: html_to_send = html_to_send[:100000] + "...(truncated)"
    prompt = f"""Analyze Moodle HTML. Find due dates/deadlines/time remaining. Extract 'status' & 'time'. If none, return JSON: {{"status": "Not Found", "time": null}}. Return ONLY valid JSON. HTML: {html_to_send}"""
    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = ai_client.generate_content(prompt)
            data = json.loads(response.text)
            if data.get("status") == "Not Found": return None
            print(f"         [AI Fallback] Extracted deadline: {data}")
            return data
        except Exception as e:
            error_str = str(e); is_rate_limit = "429" in error_str
            wait_time = 60
            if is_rate_limit:
                match = re.search(r'(?:retry(?:_delay)?|Please retry in)\s*(?:{\s*seconds:\s*|\s*)(\d+)', error_str, re.IGNORECASE)
                if match: wait_time = int(match.group(1)) + 2
                print(f"         [AI Fallback] Rate Limit (429). Wait {wait_time}s (Attempt {attempt+1}/{max_retries})...")
            else: print(f"         [AI Fallback] Failed (Attempt {attempt+1}/{max_retries}): {e}")
            if attempt + 1 < max_retries and is_rate_limit: time.sleep(wait_time)
            elif attempt + 1 >= max_retries: print("         [AI Fallback] Max retries reached."); return None
    return None

def get_deadline_info(html_content: str) -> dict | None:
    # ... (Paste your full function code here) ...
    try: soup = BeautifulSoup(html_content, 'lxml')
    except ImportError: soup = BeautifulSoup(html_content, 'html.parser')
    except Exception: return None
    data = extract_deadline_with_selectors(soup)
    if data: data["method"] = "selector"; return data
    data = extract_deadline_with_ai(html_content)
    if data: data["method"] = "ai"; return data
    return None

# (Paste the email notification helper)
import smtplib, ssl
from email.message import EmailMessage
from config import GMAIL_SENDER, GMAIL_APP_PASSWORD, GMAIL_RECEIVER
def send_email_notification(subject, body):
    """Sends an email using Gmail credentials from .env."""
    if not GMAIL_SENDER or not GMAIL_APP_PASSWORD or not GMAIL_RECEIVER:
        print("   [Email] ⚠️ Gmail credentials (SENDER, APP_PASSWORD, RECEIVER) not set. Skipping notification.")
        return

    print(f"   [Email] Connecting to Gmail to send notification to {GMAIL_RECEIVER}...")

    try:
        # Create the email message
        msg = EmailMessage()
        msg['Subject'] = subject
        msg['From'] = GMAIL_SENDER
        msg['To'] = GMAIL_RECEIVER
        msg.set_content(body)

        # Connect and send
        context = ssl.create_default_context()
        with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as smtp:
            smtp.login(GMAIL_SENDER, GMAIL_APP_PASSWORD)
            smtp.send_message(msg)
        print("   [Email] ✅ Notification email sent successfully.")
    except Exception as e:
        print(f"   [Email] ❌ Failed to send email: {e}")


# --- Main Scrape Function ---
def perform_full_scrape(lms_user, lms_pass):
    """The main scraping process, designed to run in a background thread."""
    if state.IS_SCRAPING:
        print("Scrape already in progress.")
        return
    
    state.IS_SCRAPING = True
    print(f"\n🚀 Starting full scrape for user {lms_user}...")
    
    driver = None
    db = None 
    cursor = None
    index_writer = None
    # ----------------------------------------
    
    # --- State Tracking Variables ---
    print("   [State] Loading old scrape state...")
    old_state = {"deadlines": [], "files": []}
    try:
        if os.path.exists(STATE_FILE):
            with open(STATE_FILE, 'r', encoding='utf-8') as f:
                old_state = json.load(f)
        else:
             print("   [State] No old state file found, will treat all findings as new.")
    except Exception as e:
        print(f"   [State] ⚠️ Could not load old state file: {e}")
    
    all_found_deadline_urls = set()
    all_found_file_names = set()
    # ---------------------------------

    try:
        # --- 1. Create a NEW DB connection for THIS thread ---
        print("   [DB] Scrape thread connecting to database...")
        db = sqlite3.connect(DATABASE_FILE, detect_types=sqlite3.PARSE_DECLTYPES, timeout=10)
        db.row_factory = sqlite3.Row
        cursor = db.cursor()
        print("   [DB] Scrape thread connected.")
        # -----------------------------------------------------

        # --- 2. Setup Selenium ---
        print("   Setting up WebDriver...")
        options = webdriver.ChromeOptions()
        options.add_argument("--headless")
        options.add_argument("--disable-gpu")
        options.add_argument("--window-size=1920,1080")
        options.add_argument("--no-sandbox")
        options.add_argument("--disable-dev-shm-usage")
        options.add_argument(f'user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/100.0.0.0 Safari/537.36')
        driver = webdriver.Chrome(options=options)
        driver.implicitly_wait(10)
        print("   WebDriver initialized (headless).")

        # --- 3. Login and Confirm ---
        print(f"   Attempting login...")
        driver.get("https://lms.fit.hanu.vn/login/index.php")
        try:
            username_field = WebDriverWait(driver, 15).until(EC.presence_of_element_located((By.ID, "username")))
            username_field.send_keys(lms_user)
            password_field = driver.find_element(By.ID, "password")
            login_button = driver.find_element(By.ID, "loginbtn")
            password_field.send_keys(lms_pass)
            login_button.click()
            print("   Login form submitted.")
            print("   Waiting for login result...")
            WebDriverWait(driver, 25).until(
                EC.any_of(
                    EC.presence_of_element_located((By.ID, "loggedin-user")),
                    EC.presence_of_element_located((By.CSS_SELECTOR, "div.alert.alert-danger"))
                )
            )
            
            # Check which element appeared
            try:
                driver.find_element(By.ID, "loggedin-user")
                print("   ✅ Login successful (Found 'loggedin-user' ID).")

                # --- 4. Extract Sesskey & Cookies ---
                print("      Extracting session key & cookies...")
                driver.get("https://lms.fit.hanu.vn/my/courses.php")
                time.sleep(2)
                page_source = driver.page_source

                sesskey = None
                m = re.search(r'M\.cfg\s*=\s*(\{.*?\});', page_source, re.DOTALL)
                if m:
                    cfg_json = m.group(1)
                    try: sesskey = json.loads(cfg_json).get("sesskey")
                    except json.JSONDecodeError:
                         m2 = re.search(r'"sesskey"\s*:\s*"([^"]+)"', cfg_json); sesskey = m2.group(1) if m2 else None
                if not sesskey:
                    m3 = re.search(r'"sesskey"\s*:\s*"([^"]+)"', page_source); sesskey = m3.group(1) if m3 else None

                if not sesskey:
                    print("      ❌ Failed to extract sesskey from /my/courses.php.")
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    try:
                         driver.save_screenshot(f"sesskey_fail_screenshot_{timestamp}.png")
                         with open(f"sesskey_fail_source_{timestamp}.html", "w", encoding="utf-8") as f: f.write(page_source)
                         print(f"         📸📄 Saved debug info for sesskey failure.")
                    except Exception as save_e: print(f"         ⚠️ Could not save debug info: {save_e}")
                    raise ValueError("Sesskey not found post-login")

                print(f"      🔑 Sesskey: {sesskey}")
                
                selenium_cookies = driver.get_cookies()
                cookies_dict = {c['name']: c['value'] for c in selenium_cookies}
                headers = {"accept": "application/json, text/javascript, */*; q=0.01", "content-type": "application/json",
                           "origin": "https://lms.fit.hanu.vn", "referer": "https://lms.fit.hanu.vn/my/courses.php",
                           "user-agent": driver.execute_script("return navigator.userAgent;"), "x-requested-with": "XMLHttpRequest",}

            except NoSuchElementException: # Login failed
                error_message = "Unknown login error"
                try: error_message = driver.find_element(By.CSS_SELECTOR, "div.alert.alert-danger").text
                except NoSuchElementException: pass
                print(f"   ❌ Login failed: {error_message}")
                raise Exception(f"Login failed: {error_message}")

        except TimeoutException:
            print("   ❌ Login failed: Timeout waiting for success/failure element.")
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S"); driver.save_screenshot(f"login_timeout_{timestamp}.png")
            raise
        except Exception as login_e:
            print(f"   ❌ Login failed: {login_e}")
            raise

        # --- 5. Fetch Course List (AJAX) ---
        print("   Fetching course list...")
        ajax_url = f"https://lms.fit.hanu.vn/lib/ajax/service.php?sesskey={sesskey}&info=core_course_get_enrolled_courses_by_timeline_classification"
        payload = [{"index": 0, "methodname": "core_course_get_enrolled_courses_by_timeline_classification",
                    "args": {"offset": 0, "limit": 999, "classification": "all", "sort": "fullname"}}]
        response = requests.post(ajax_url, headers=headers, cookies=cookies_dict, json=payload, timeout=REQUESTS_TIMEOUT)
        response.raise_for_status(); json_data = response.json()
        simplified_courses = []
        if json_data and isinstance(json_data[0], dict) and json_data[0].get("error") is False:
            simplified_courses = [{"id": c.get("id"), "name": c.get("fullname"), "url": c.get("viewurl")}
                                  for c in json_data[0].get("data", {}).get("courses", [])]
        else: print(f"   ❌ AJAX error: {json_data}"); raise ValueError("Course fetch failed")

        # --- 6. [DB] Clear old data and insert new courses ---
        print("   [DB] Clearing old courses and deadlines...")
        cursor.execute('DELETE FROM deadlines')
        cursor.execute('DELETE FROM courses')
        print(f"   [DB] Inserting {len(simplified_courses)} new courses...")
        for course in simplified_courses:
             cursor.execute(
                 'INSERT INTO courses (course_id, name, url) VALUES (?, ?, ?)',
                 (course.get('id'), course.get('name'), course.get('url'))
             )
        print(f"   📝 Saved {len(simplified_courses)} courses to database.")
        # --- [END DB] ---

        # --- 7. [Whoosh] Prepare Search Index ---
        print("   [Search] Clearing and opening index writer...")
        search_index = clear_search_index()
        index_writer = search_index.writer()
        # --- [END Whoosh] ---

        # --- 8. Course Processing Loop ---
        for course in simplified_courses:
            course_id = course.get("id"); course_name = course.get("name", "N/A"); course_url = course.get("url")
            safe_course_name = re.sub(r'[\\/*?:"<>|]', "_", course_name).strip()[:150]
            course_folder = os.path.join(SAVE_DIR, f"{course_id}_{safe_course_name}")
            os.makedirs(course_folder, exist_ok=True)
            print(f"\n   📘 Processing course {course_id} - {course_name}")
            if not course_url: print("      ⚠️ No URL, skipping."); continue

            try:
                driver.get(course_url)
                WebDriverWait(driver, 20).until(EC.presence_of_element_located((By.ID, "region-main")))
            except Exception as e: print(f"      ⚠️ Failed load main page: {e}"); continue

            main_page_source = driver.page_source
            main_filename = os.path.join(course_folder, "main_page.html")
            with open(main_filename, "w", encoding="utf-8") as fh: fh.write(main_page_source)
            print(f"      💾 Saved main HTML -> {main_filename}")

            # Collect sub-links
            links_to_visit = []; visited_hrefs = set(); parsed_course_netloc = urllib.parse.urlparse(course_url).netloc
            try: anchor_elements = driver.find_elements(By.CSS_SELECTOR, "#region-main a[href]")
            except Exception: anchor_elements = driver.find_elements(By.CSS_SELECTOR, "a[href]")

            for a in anchor_elements:
                try:
                    href_raw = a.get_attribute("href")
                    if not href_raw: continue
                    href = urllib.parse.urljoin(course_url, href_raw).split("#")[0]
                    if not href or href in visited_hrefs or href == course_url: continue
                    if href.lower().startswith(("javascript:", "mailto:", "tel:")): continue
                    parsed_href_netloc = urllib.parse.urlparse(href).netloc
                    if parsed_href_netloc and parsed_href_netloc != parsed_course_netloc: continue
                    link_text = a.text.strip()
                    if not link_text:
                        try: fn_span = a.find_element(By.CSS_SELECTOR, "span.fp-filename"); link_text = fn_span.text.strip()
                        except: link_text = ""
                    links_to_visit.append((href, link_text)); visited_hrefs.add(href)
                except Exception: continue

            total_to_visit = len(links_to_visit) if (MAX_SUBPAGES is None) else min(len(links_to_visit), MAX_SUBPAGES)
            print(f"      🔗 Found {len(links_to_visit)} links, processing {total_to_visit}...")
            if total_to_visit == 0: continue

            deadlines_to_add = [] # List to hold deadlines for this course

            # --- 9. Visit Subpages Loop ---
            for idx, (href, link_text) in enumerate(links_to_visit[:total_to_visit], start=1):
                print(f"\n         👉 [{idx}/{total_to_visit}] Visiting: {href}")
                if link_text: print(f"            Link Text: '{link_text}'")

                file_extensions = (".pdf", ".docx", ".pptx", ".zip", ".rar", ".xls", ".xlsx")
                is_direct_file = "pluginfile.php" in href or \
                                 any(href.lower().endswith(ext) for ext in file_extensions) or \
                                 (link_text and any(link_text.lower().endswith(ext) for ext in file_extensions))

                if is_direct_file:
                    local_path = download_file(href, course_folder, cookies_dict, headers, link_text)
                    if local_path:
                        # --- [STATE] Track this file ---
                        all_found_file_names.add(os.path.basename(local_path))
                        # -----------------------------
                        extracted_text=None; file_type="Unknown"; file_ext_lower=os.path.splitext(local_path)[1].lower()
                        if file_ext_lower==".docx": file_type="Word"; extracted_text=read_docx(local_path)
                        elif file_ext_lower==".pptx": file_type="PowerPoint"; extracted_text=read_pptx(local_path)
                        elif file_ext_lower==".pdf": file_type="PDF"; extracted_text=read_pdf(local_path)
                        
                        if extracted_text:
                            cleaned_text = clean_file_text(extracted_text)
                            txt_fname = f"{os.path.splitext(os.path.basename(local_path))[0]}.txt"
                            txt_fpath = os.path.join(course_folder, txt_fname)
                            with open(txt_fpath, "w", encoding="utf-8") as f: f.write(cleaned_text)
                            print(f"            💾 Saved cleaned text -> {txt_fpath}")

                            print(f"            [Search] Indexing {os.path.basename(local_path)}...")
                            index_writer.add_document(
                                course_id=str(course_id), course_name=course_name,
                                file_name=os.path.basename(local_path),
                                file_type=file_type, content=cleaned_text
                            )
                        elif not extracted_text: print("            ➖ Skipping (no text/unsupported/archive).")
                    continue # Next link in main list

                # --- Process as HTML Page ---
                else:
                    try:
                        print("         📄 Visiting as HTML page...")
                        driver.get(href)
                        WebDriverWait(driver, 15).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
                        time.sleep(1)
                        current_page_source = driver.page_source
                    except Exception as e:
                        print(f"            ⚠️ Failed load page: {e}"); continue

                    # Save HTML
                    parsed=urllib.parse.urlparse(href); path_part=parsed.path.strip("/") or "root"; query_part=parsed.query.replace("&","_").replace("=","-")
                    safe_subname=re.sub(r'[\\/*?:"<>|]', "_", f"{idx}_{path_part}" + (f"_{query_part}" if query_part else ""))[:200]
                    sub_filename=os.path.join(course_folder, f"{safe_subname}.html")
                    try:
                        with open(sub_filename, "w", encoding="utf-8") as fh: fh.write(current_page_source)
                        print(f"            💾 Saved HTML -> {sub_filename}")
                    except Exception as save_html_e: print(f"            ⚠️ Failed to save HTML: {save_html_e}")

                    # Find and Process Nested File Links
                    try:
                        print("         DEBUG: Waiting up to 10s for nested pluginfile links...")
                        try:
                            WebDriverWait(driver, 10).until(
                                EC.presence_of_element_located((By.CSS_SELECTOR, "a[href*='pluginfile.php']"))
                            )
                            nested_links = driver.find_elements(By.CSS_SELECTOR, "a[href*='pluginfile.php']")
                            print(f"         DEBUG: Found {len(nested_links)} potential nested links after wait.")
                        except TimeoutException:
                            print("         ➖ No nested pluginfile links found within 10 seconds.")
                            nested_links = []

                        unique_nested_hrefs = set()
                        processed_nested_count = 0

                        if nested_links:
                            print(f"         🔎 Processing {len(nested_links)} potential nested file link(s)...")
                            for file_link_element in nested_links:
                                 try:
                                     f_href_raw = file_link_element.get_attribute('href')
                                     f_href = urllib.parse.urljoin(href, f_href_raw).split("#")[0]
                                     if not f_href or f_href in unique_nested_hrefs: continue
                                     unique_nested_hrefs.add(f_href); processed_nested_count += 1
                                     
                                     f_link_text = file_link_element.text.strip()
                                     if not f_link_text:
                                          try: f_link_text=file_link_element.find_element(By.CSS_SELECTOR, "span.fp-filename").text.strip()
                                          except: f_link_text=os.path.basename(urllib.parse.urlparse(f_href).path)

                                     print(f"            📂 Processing nested file [{processed_nested_count}]: {f_href} (Text: '{f_link_text}')")
                                     nested_local_path = download_file(f_href, course_folder, cookies_dict, headers, f_link_text)
                                     
                                     if nested_local_path:
                                          # --- [STATE] Track this file ---
                                          all_found_file_names.add(os.path.basename(nested_local_path))
                                          # -----------------------------
                                          extracted_text=None; file_type="Unknown"; file_ext_lower=os.path.splitext(nested_local_path)[1].lower()
                                          if file_ext_lower==".docx": file_type="Word"; extracted_text=read_docx(nested_local_path)
                                          elif file_ext_lower==".pptx": file_type="PowerPoint"; extracted_text=read_pptx(nested_local_path)
                                          elif file_ext_lower==".pdf": file_type="PDF"; extracted_text=read_pdf(nested_local_path)

                                          if extracted_text:
                                               cleaned_text = clean_file_text(extracted_text)
                                               txt_fname=f"{os.path.splitext(os.path.basename(nested_local_path))[0]}.txt"
                                               txt_fpath=os.path.join(course_folder, txt_fname)
                                               with open(txt_fpath, "w", encoding="utf-8") as f: f.write(cleaned_text)
                                               print(f"               💾 Saved cleaned text -> {txt_fpath}")

                                               print(f"               [Search] Indexing {os.path.basename(nested_local_path)}...")
                                               index_writer.add_document(
                                                   course_id=str(course_id), course_name=course_name,
                                                   file_name=os.path.basename(nested_local_path),
                                                   file_type=file_type, content=cleaned_text
                                               )
                                          elif not extracted_text: print("               ➖ Skipping nested analysis (no text/unsupported/archive).")
                                     else: print(f"            ❌ Nested download FAILED for {f_href}.")
                                 except Exception as nested_proc_e:
                                     print(f"               ⚠️ Error processing nested link {f_href or 'unknown'}: {nested_proc_e}")
                            
                            if processed_nested_count == 0 and nested_links:
                                 print("         ➖ No unique/valid nested pluginfile links after filtering.")
                        
                    except Exception as sub_e:
                        print(f"            ⚠️ Error finding/processing nested links: {sub_e}")
                    # --- End Nested File Link Processing ---

                    # --- Check Deadlines ---
                    if "mod/assign/" in href or "mod/quiz/" in href:
                        deadline_info = get_deadline_info(current_page_source)
                        if deadline_info:
                            # --- [STATE] Track this deadline ---
                            all_found_deadline_urls.add(href) 
                            # -----------------------------------
                            
                            deadline_info["url"] = href
                            original_time_str = deadline_info.get("time")
                            iso_timestamp = None
                            
                            if original_time_str and deadline_info.get("status") == "Due":
                                try:
                                    parsed_datetime = date_parser.parse(original_time_str, fuzzy=True)
                                    iso_timestamp = parsed_datetime.isoformat()
                                except Exception as parse_e:
                                    print(f"               ⚠️ Could not parse deadline string '{original_time_str}': {parse_e}")
                            
                            elif original_time_str and deadline_info.get("status") == "Time Remaining":
                                time_delta = parse_time_remaining(original_time_str)
                                if time_delta:
                                    now_local = datetime.now()
                                    calculated_due_date = now_local + time_delta
                                    iso_timestamp = calculated_due_date.isoformat()
                                else:
                                    print(f"               ⚠️ Could not calculate relative time: '{original_time_str}'")

                            deadlines_to_add.append((
                                course_id, deadline_info.get('status'),
                                deadline_info.get('time'), iso_timestamp, href
                            ))
                            print(f"            🎯 Deadline Found (Method: {deadline_info.get('method', 'N/A')}): {deadline_info.get('status','N/A')} - {original_time_str}")
            # --- End Subpage Loop ---

            # --- Save deadlines for this course to DB ---
            if deadlines_to_add:
                print(f"      [DB] Inserting {len(deadlines_to_add)} deadlines for course {course_id}...")
                cursor.executemany(
                    'INSERT INTO deadlines (course_id, status, time_string, parsed_iso_date, url) VALUES (?, ?, ?, ?, ?)',
                    deadlines_to_add
                )
        # --- End Course Loop ---

        # --- 10. Commit all changes ---
        print("\n   [Search] Committing index writer...")
        index_writer.commit()
        index_writer = None # Mark as committed
        print("   [Search] Index commit complete.")

        print("\n   [DB] Committing all scrape data to database...")
        db.commit() # Commit all DB inserts/deletes
        
        # --- [NEW] Compare State and Send Notification ---
        print("\n   [State] Comparing scrape results to previous state...")
        old_deadline_set = set(old_state.get("deadlines", []))
        old_file_set = set(old_state.get("files", []))

        # Find items in the new sets that were not in the old sets
        new_deadlines = all_found_deadline_urls - old_deadline_set
        new_files = all_found_file_names - old_file_set
        
        email_body_lines = []
        email_subject = "LMS Assistant: No New Updates"

        if new_deadlines:
            print(f"   [State] ✅ Found {len(new_deadlines)} new/updated deadline(s)!")
            email_body_lines.append(f"Found {len(new_deadlines)} new or changed deadline pages:")
            for url in new_deadlines:
                email_body_lines.append(f"- {url}")
        
        if new_files:
            print(f"   [State] ✅ Found {len(new_files)} new file(s)!")
            email_body_lines.append(f"\nFound {len(new_files)} new file(s):")
            for name in new_files:
                email_body_lines.append(f"- {name}")
        
        if not new_deadlines and not new_files:
            print("   [State] ➖ No new deadlines or files found.")
        
        # Send email only if there are updates
        if email_body_lines:
            email_subject = f"LMS Assistant: {len(new_deadlines)} New Deadlines, {len(new_files)} New Files!"
            final_email_body = "Your LMS Assistant scrape found the following updates:\n\n" + "\n".join(email_body_lines)
            send_email_notification(email_subject, final_email_body)
        
        # Save the *current* state for next time
        print("   [State] Saving current state...")
        new_state_data = {
            "deadlines": list(all_found_deadline_urls),
            "files": list(all_found_file_names)
        }
        with open(STATE_FILE, 'w', encoding='utf-8') as f:
            json.dump(new_state_data, f, indent=4)
        # --- [END NEW] ---
        
        print("\n✅ Full scrape completed successfully.")

    except Exception as scrape_e:
        print(f"\n❌ An error occurred during scraping: {scrape_e}")
        import traceback
        traceback.print_exc()
        if db:
             print("   [DB] Rolling back database changes due to error.")
             db.rollback()
        if index_writer:
             print("   [Search] Cancelling index writer due to error.")
             index_writer.cancel()
    finally:
        if db:
             print("   [DB] Closing database connection...")
             db.close()
        if driver:
            print("   Closing WebDriver...")
            driver.quit()
        state.IS_SCRAPING = False # Reset flag
        print("   Scrape function finished.")